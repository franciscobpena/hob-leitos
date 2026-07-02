# -*- coding: utf-8 -*-
# Gera o deck "Projetos Kaizen A3 — HOB Visita 06" a partir do hob_a3_data.json
import json, os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

BASE = os.path.dirname(os.path.abspath(__file__))
data = json.load(open(os.path.join(BASE, 'hob_a3_data.json'), encoding='utf-8'))
IND = data['indicadores']

NAVY   = RGBColor(0x0B, 0x1A, 0x33)
NAVY2  = RGBColor(0x13, 0x2B, 0x4F)
TEAL   = RGBColor(0x0D, 0x94, 0x88)
GREEN  = RGBColor(0x1F, 0x8A, 0x4C)
AMBER  = RGBColor(0xB4, 0x69, 0x00)
RED    = RGBColor(0xB0, 0x20, 0x20)
GRAY   = RGBColor(0x66, 0x70, 0x85)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
DARK   = RGBColor(0x1A, 0x1A, 0x1A)
LIGHT  = RGBColor(0xF2, 0xF5, 0xF9)
BORDER = RGBColor(0xB9, 0xC6, 0xCF)

STATUS = {
    'concluido':            ('Concluído', GREEN),
    'em_andamento':         ('Em andamento', AMBER),
    'em_atraso':            ('Em atraso', RED),
    'estacionamento_ideias':('Estac. de ideias', GRAY),
}

RAIA_COLOR = {'Governança': TEAL, 'Passagem': RGBColor(0x2F,0x6F,0xB7), 'Saída': RGBColor(0x7A,0x4F,0xB0)}

def fmt_prazo(v):
    if not v: return '—'
    y, m, d = v.split('-')
    return f'{d}/{m}/{y}'

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]

def txbox(slide, x, y, w, h, text, size, color=DARK, bold=False, align=PP_ALIGN.LEFT, font='Segoe UI', anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    lines = text.split('\n')
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run(); r.text = ln
        r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = color; r.font.name = font
    return tb

def rect(slide, x, y, w, h, fill, line=None):
    from pptx.enum.shapes import MSO_SHAPE
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if line: sh.line.color.rgb = line; sh.line.width = Pt(0.75)
    else: sh.line.fill.background()
    sh.shadow.inherit = False
    return sh

# ─── CAPA ───
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, prs.slide_width, prs.slide_height, NAVY)
rect(s, 0, Inches(3.42), prs.slide_width, Pt(3), TEAL)
txbox(s, Inches(0.9), Inches(1.5), Inches(11.5), Inches(0.6), 'PROJETO LEAN NAS EMERGÊNCIAS', 20, TEAL, True)
txbox(s, Inches(0.9), Inches(2.0), Inches(11.5), Inches(1.2), 'Projetos Kaizen (A3) — HOB', 40, WHITE, True)
txbox(s, Inches(0.9), Inches(3.6), Inches(11.5), Inches(0.5), 'Hospital Metropolitano Odilon Behrens · Belo Horizonte/MG', 16, RGBColor(0xB8,0xC4,0xD6))
txbox(s, Inches(0.9), Inches(4.15), Inches(11.5), Inches(0.5), 'Visita 06 · 06 e 07/07/2026', 16, WHITE, True)
txbox(s, Inches(0.9), Inches(5.3), Inches(11.5), Inches(0.9),
      'Premissa da apresentação: preenchimento dos A3 e validação das metas SMART e indicadores.\nPlanos de ação atualizados na plataforma: hob-leitos.vercel.app/#/kaizens', 13, RGBColor(0xB8,0xC4,0xD6))
txbox(s, Inches(0.9), Inches(6.8), Inches(11.5), Inches(0.4), 'BP | Sírio-Libanês · CONASS · CONASEMS · PROADI-SUS · SUS · Ministério da Saúde', 10, RGBColor(0x8A,0x9A,0xB5))

# ─── GUIA ───
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, prs.slide_width, Inches(0.9), NAVY)
txbox(s, Inches(0.5), Inches(0.18), Inches(12), Inches(0.55), 'Como apresentar cada A3 — 15 minutos por frente', 20, WHITE, True)
guia = [
    ('1. Meta SMART e indicador', 'Posicionar a meta (SMART) e o indicador do projeto, já no modelo A3. Validar com o grupo se a meta está mensurável e com prazo.'),
    ('2. Resultados', 'Mostrar o comportamento do indicador do projeto: onde estava, onde está, distância da meta.'),
    ('3. Principais ações', 'O que foi concluído desde a última visita e o que está em andamento (com farol de status).'),
    ('4. Bloqueios e próximos passos', 'O que trava o avanço, o que precisa de decisão da Diretoria e quais os próximos passos com prazo e responsável.'),
]
y = 1.3
for t, d in guia:
    txbox(s, Inches(0.7), Inches(y), Inches(4.0), Inches(0.5), t, 16, TEAL, True)
    txbox(s, Inches(4.9), Inches(y), Inches(7.9), Inches(0.9), d, 13, DARK)
    y += 1.15
txbox(s, Inches(0.7), Inches(6.2), Inches(12), Inches(0.6),
      'Antes da visita: atualizar o plano de ação do seu projeto na plataforma — hob-leitos.vercel.app/#/kaizens', 13, NAVY2, True)

# ─── SLIDES POR PROJETO ───
for p in data['projetos']:
    s = prs.slides.add_slide(BLANK)
    # header
    rect(s, 0, 0, prs.slide_width, Inches(0.95), NAVY)
    txbox(s, Inches(0.4), Inches(0.10), Inches(10.2), Inches(0.75), p['nome'], 17, WHITE, True)
    rc = RAIA_COLOR.get(p['raia'], TEAL)
    rect(s, Inches(11.0), Inches(0.24), Inches(1.95), Inches(0.45), rc)
    txbox(s, Inches(11.0), Inches(0.27), Inches(1.95), Inches(0.4), p['raia'], 12, WHITE, True, PP_ALIGN.CENTER)
    # linha de responsáveis
    resp = f"Responsável: {p['responsavel']}"
    if p.get('corresponsavel'): resp += f"  ·  Corresponsável: {p['corresponsavel']}"
    resp += f"  ·  Prazo do projeto: {fmt_prazo(p.get('prazo'))}"
    txbox(s, Inches(0.4), Inches(1.02), Inches(12.5), Inches(0.35), resp, 12, GRAY, True)
    # meta smart
    rect(s, Inches(0.4), Inches(1.45), Inches(12.53), Inches(0.85), LIGHT, TEAL)
    txbox(s, Inches(0.55), Inches(1.51), Inches(1.6), Inches(0.4), 'META SMART', 11, TEAL, True)
    txbox(s, Inches(2.15), Inches(1.49), Inches(10.6), Inches(0.8), p['meta_smart'], 11.5, DARK)
    # indicadores
    inds = ' · '.join(IND.get(i, i) for i in p['indicadores'])
    txbox(s, Inches(0.4), Inches(2.38), Inches(12.5), Inches(0.32), f'Indicadores vinculados: {inds}', 10.5, NAVY2)
    # tabela de ações
    acoes = p['acoes'][:7]
    extra = len(p['acoes']) - len(acoes)
    rows = len(acoes) + 1
    tbl_h = Inches(0.34 * rows)
    gfx = s.shapes.add_table(rows, 4, Inches(0.4), Inches(2.80), Inches(12.53), tbl_h)
    tbl = gfx.table
    tbl.columns[0].width = Inches(7.6)
    tbl.columns[1].width = Inches(2.0)
    tbl.columns[2].width = Inches(1.15)
    tbl.columns[3].width = Inches(1.78)
    heads = ['Ação', 'Responsável', 'Prazo', 'Status']
    for j, htxt in enumerate(heads):
        c = tbl.cell(0, j)
        c.fill.solid(); c.fill.fore_color.rgb = NAVY
        c.margin_top = Pt(2); c.margin_bottom = Pt(2)
        tf = c.text_frame; tf.word_wrap = True
        pr = tf.paragraphs[0]; r = pr.add_run(); r.text = htxt
        r.font.size = Pt(10.5); r.font.bold = True; r.font.color.rgb = WHITE; r.font.name = 'Segoe UI'
    for i, a in enumerate(acoes, start=1):
        stx, scol = STATUS.get(a['status'], (a['status'], GRAY))
        pct = a.get('pct')
        if a['status'] == 'em_andamento' and pct: stx += f' ({pct}%)'
        vals = [a['descricao'], a.get('responsavel') or '—', fmt_prazo(a.get('prazo')), stx]
        for j, v in enumerate(vals):
            c = tbl.cell(i, j)
            c.fill.solid(); c.fill.fore_color.rgb = WHITE if i % 2 else LIGHT
            c.margin_top = Pt(1); c.margin_bottom = Pt(1)
            tf = c.text_frame; tf.word_wrap = True
            pr = tf.paragraphs[0]; r = pr.add_run(); r.text = str(v)
            r.font.size = Pt(9.5); r.font.name = 'Segoe UI'
            r.font.color.rgb = scol if j == 3 else DARK
            r.font.bold = (j == 3)
    foot = 'HOB · Visita 06 · plataforma: hob-leitos.vercel.app/#/kaizens'
    if extra > 0: foot = f'+{extra} outra(s) ação(ões) na plataforma  ·  ' + foot
    txbox(s, Inches(0.4), Inches(7.08), Inches(12.5), Inches(0.32), foot, 9.5, GRAY)

out = os.path.join(BASE, 'HOB-Projetos-Kaizen-A3-Visita06.pptx')
prs.save(out)
print('OK', out)
